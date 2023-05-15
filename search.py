import os
from sys import exit
from openpyxl import Workbook
import comtypes.client
from pptx import Presentation
# https://pymupdf.readthedocs.io/en/latest/module.html
# https://pymupdf.readthedocs.io/en/latest/document.html#Document.needs_pass
import fitz 
# https://docs.aspose.com/slides/python-net/password-protected-presentation/#checking-whether-a-presentation-is-password-protected-before-loading-it
from aspose import slides

def app_quit():
    input("\nPress enter to close")
    exit()

def key_quit():
    print("Program canceled by user")
    exit()

def search_files(search_path, keywords, workbook_xlsx, workbook_path):
    try:
        # Create a Word Application object
        word_app = comtypes.client.CreateObject("Word.Application")
        # Create a PowerPoint Application objects
        ppt_app = comtypes.client.CreateObject("PowerPoint.Application")
        # ppt_app.Visible = True
    except Exception as e:
        print(f"Error creating objects: {e}")
        exit()

    # Create new workbook
    workbook = Workbook()
    workbook_name = workbook_xlsx
    workbook_path = workbook_path + "/" + workbook_name
    worksheet = workbook.active
    worksheet.title = "File List"
    worksheet.column_dimensions['A'].width = 30
    worksheet.column_dimensions['B'].width = 120
    worksheet.column_dimensions['C'].width = 30
    row = 2
    # worksheet header information
    worksheet.cell(row=1, column=1).value = "File Name"
    worksheet.cell(row=1, column=2).value = "File Path"
    worksheet.cell(row=1, column=3).value = "Keywords Found"

    # wrong password to skip password required files
    pw = "Password"
    # limit search to large file size to reduce long processing times (100 KB)
    # MAX_FILE_SIZE = 100000000
    try:
        for root, dirs, files in os.walk(search_path):
            for filename in files:
                if filename.endswith(".doc") or filename.endswith(".docx"):
                    try:
                        doc_path = os.path.join(root, filename)
                        print(f"Opening {doc_path}")
                        try:
                            doc = word_app.Documents.Open(doc_path, ReadOnly=True, PasswordDocument=pw)
                        except:
                            print(f"Password protected: Skipping {doc_path}")
                            continue
                        text = doc.Content.Text.lower()
                        # text = text.replace('-', ' ')

                        found_words = set()
                        # Check if any of the search words are in the text
                        for phrase in keywords:
                            if phrase in text:
                                found_words.add(phrase)
                        # Write the file information and found words to the Excel worksheet
                        if found_words:
                            print(f"Found {found_words} in {filename}")
                            worksheet.cell(row=row, column=1).value = filename
                            worksheet.cell(row=row, column=2).value = doc_path
                            worksheet.cell(row=row, column=2).hyperlink = doc_path
                            worksheet.cell(row=row, column=3).value = ", ".join(found_words)
                            row += 1
                            
                    except Exception as e:
                        print(f"Error reading {doc_path}: {e}")
                    
                    finally:
                        try:
                            print(f"Closing {doc_path}")
                            doc.Close(SaveChanges=False)
                        except:
                            print(f"Error closing {doc_path}")
                    
                elif filename.endswith(".ppt") or filename.endswith(".pptx"):
                    ppt_path = os.path.join(root, filename)
                    # skip large files
                    # try:
                    #     if os.path.getsize(ppt_path) > MAX_FILE_SIZE:
                    #         print(f"Skipping {ppt_path} because it is too large")
                    #         continue
                    # except Exception as e:
                    #     print(f"Error getting file size for {filename}: {e}")

                    # Open the PowerPoint file
                    try:
                        # skip if password protected
                        print(f"Opening {ppt_path}")
                        # https://docs.aspose.com/slides/python-net/password-protected-presentation/
                        ppt = slides.PresentationFactory.instance.get_presentation_info(ppt_path)
                        if ppt.is_password_protected:
                            print(f"Password protected: Skipping {ppt_path}")
                            continue
                        found_words = set()
                        # don't need to close when using python-pptx
                        if filename.endswith(".pptx"):
                            ppt = Presentation(ppt_path)
                                
                            # Iterate through every slide
                            for slide in ppt.slides:
                                # Iterate through all shapes with text attributes
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text = shape.text.lower()
                                        # Check if any of the search words exist in the shape's text
                                        for phrase in keywords:
                                            if phrase in text:
                                                found_words.add(phrase)
                        # work around for ppt files
                        else:
                            ppt = ppt_app.Presentations.Open(ppt_path, ReadOnly=True)
                            # Iterate through every slide
                            for slide in ppt.Slides:
                                # Iterate through all the shapes in the slide
                                for shape in slide.Shapes:
                                    # Type 6 shape causing errors (can't get range of text in group of shapes)
                                    if shape.Type == 6:
                                        continue
                                    # Check if the shape has text
                                    if hasattr(shape, "TextFrame"):
                                        # Get the text from the shape
                                        if shape.TextFrame.HasText:
                                            text = shape.TextFrame.TextRange.Text.lower()
                                            # Check if any of the search words exist in the shape's text
                                            for phrase in keywords:
                                                if phrase in text:
                                                    found_words.add(phrase)
                            
                            
                        if found_words:
                            print(f"Found {found_words} in {filename}")
                            worksheet.cell(row=row, column=1).value = filename
                            worksheet.cell(row=row, column=2).value = ppt_path
                            worksheet.cell(row=row, column=2).hyperlink = ppt_path
                            worksheet.cell(row=row, column=3).value = ", ".join(found_words)
                            row += 1
 
                    except Exception as e:
                        print(f"Error reading {ppt_path}: {e}")
                    
                    finally:
                        # Close the PowerPoint file
                        if filename.endswith(".ppt"):
                            try:
                                print(f"Closing {ppt_path}")
                                ppt.Close()
                            except:
                                print(f"Error closing {ppt_path}")

                elif filename.endswith(".pdf"):
                    pdf_path = os.path.join(root, filename)
                    # skip large files
                    # try:
                    #     if os.path.getsize(pdf_path) > MAX_FILE_SIZE:
                    #         print(f"Skipping {pdf_path} because it is too large")
                    #         continue
                    # except Exception as e:
                    #     print(f"Error getting file size for {filename}: {e}")
                    try:
                        print(f"Opening {pdf_path}")
                        found_words = set()
                        pdf_file = fitz.open(pdf_path)
                        if pdf_file.needs_pass:
                            print(f"Password protected: Skipping {pdf_path}")
                            continue
                        for page in pdf_file:
                            text = page.get_text().lower()
                            # text = text.replace('-', ' ')
                            
                            # Check if any of the search words are in the text
                            for phrase in keywords:
                                if phrase in text:
                                    found_words.add(phrase)

                        if found_words:
                            print(f"Found {found_words} in {filename}")
                            worksheet.cell(row=row, column=1).value = filename
                            worksheet.cell(row=row, column=2).value = pdf_path
                            worksheet.cell(row=row, column=2).hyperlink = pdf_path
                            worksheet.cell(row=row, column=3).value = ", ".join(found_words)
                            row += 1
 

                    except Exception as e:
                        print(f"Error reading {pdf_path}: {e}")

                    finally:
                        try:
                            print(f"Closing {pdf_path}")
                            pdf_file.close()
                        except:
                            print(f"Error closing {pdf_path}")        

        if row > 2:
            print("Search phrases found")
        else:
            print("No search phrases found")
        
        try:
            workbook.save(workbook_path)
            print(f"Results saved in {workbook_path}")
            # Close the application objects
            word_app.Quit()
            ppt_app.Quit()
        finally:
            app_quit()

    except Exception as e:
        print(f"An error occured: {e}")